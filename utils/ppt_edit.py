from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

import matplotlib.pyplot as plt

from pptx.util import Inches

def insert_image_to_slide(slide, image_path, prs):
    if not os.path.exists(image_path):
        print(f"Image file not found: {image_path}")
        return

    #Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    #Set image dimensions
    image_width = Inches(9.5)
    image_height = Inches(3.5)

    #Center the image
    left = (slide_width - image_width) / 2
    top = (slide_height - image_height) / 2

    slide.shapes.add_picture(image_path, left, top, width=image_width, height=image_height)

def update_textbox_with_resource_name(slide, label_prefix, resource_name):
    for shape in slide.shapes:
        if shape.has_text_frame and label_prefix in shape.text:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = f"{label_prefix} {resource_name}"
            font = run.font
            font.name = "Inter"
            font.size = Pt(48)
            font.bold = True
            p.alignment = PP_ALIGN.LEFT
            return True
    return False


def find_slide_by_title(prs, target_title):
    target_title_lower = target_title.lower().strip()

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            try:
                if hasattr(shape, 'text') and shape.text:
                    shape_text = shape.text.lower().strip()
                    if target_title_lower in shape_text or shape_text in target_title_lower:
                        print(f"Found slide with title: '{shape.text}' (slide {slide_idx + 1})")
                        return slide

                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    if shape.placeholder_format.type == 1:
                        if hasattr(shape, 'text') and shape.text:
                            shape_text = shape.text.lower().strip()
                            if target_title_lower in shape_text or shape_text in target_title_lower:
                                print(f"Found slide with title placeholder: '{shape.text}' (slide {slide_idx + 1})")
                                return slide
            except (AttributeError, ValueError):
                continue

    print(f"Could not find slide with title containing: '{target_title}'")
    return None


def find_table_in_slide(slide):
    for shape in slide.shapes:
        if shape.has_table:
            return shape
    return None


def adjust_table_rows(table, data_count):
    current_rows = len(table.rows)
    needed_rows = data_count + 1

    if current_rows < needed_rows:
        rows_to_add = needed_rows - current_rows
        print(f"Adding {rows_to_add} rows to table")
        for _ in range(rows_to_add):
            try:
                last_row = table.rows[-1]._tr
                new_row = last_row.__class__(last_row.xml, last_row)
                table._tbl.append(new_row)
            except Exception as e:
                print(f"Error adding row: {e}")
                break

    elif current_rows > needed_rows >= 1:
        rows_to_remove = current_rows - needed_rows
        max_removable = current_rows - 1
        actual_rows_to_remove = min(rows_to_remove, max_removable)

        if actual_rows_to_remove > 0:
            print(f"Removing {actual_rows_to_remove} empty rows from table")
            try:
                for _ in range(actual_rows_to_remove):
                    if len(table.rows) > 1:
                        row_to_remove = table.rows[-1]._tr
                        table._tbl.remove(row_to_remove)
                    else:
                        break
            except Exception as e:
                print(f"Error removing rows: {e}")
def update_resource_counts_on_slide(prs, slide_index, ec2_count, rds_count, total_bill_amount=None):
    slide = prs.slides[slide_index]

    def update_shape_text(shape):
        if not shape.has_text_frame:
            return

        raw_text = shape.text_frame.text.strip()
        raw_text_lower = raw_text.lower()

        def set_text_with_format(new_text):
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            run = p.add_run()
            run.text = new_text

            # Set font
            font = run.font
            font.name = 'Arial'
            font.size = Pt(20)

            # Set alignment
            p.alignment = PP_ALIGN.CENTER

        print(f"[DEBUG] Found shape text: '{raw_text}'")

        if "number of ec2" in raw_text_lower:
            set_text_with_format(f"Number of ec2 instances: {ec2_count}")
            print("➡️  Updated EC2 count")
        elif "no. of databases" in raw_text_lower:
            set_text_with_format(f"No. Of Databases: {rds_count}")
            print("➡️  Updated RDS count")
        elif "total bill amount" in raw_text_lower and total_bill_amount is not None:
            set_text_with_format(f"Total Bill Amount: ${total_bill_amount:,.2f}")
            print("➡️  Updated Billing Amount")

    for shape in slide.shapes:
        if shape.shape_type == 6:  # Group shape
            for sub_shape in shape.shapes:
                update_shape_text(sub_shape)
        else:
            update_shape_text(shape)

    print(f"\n✅ Slide {slide_index + 1} updated with:")
    print(f"   - EC2 instances: {ec2_count}")
    print(f"   - RDS databases: {rds_count}")
    # if total_bill_amount is not None:
    #     print(f"   - Total Bill Amount: ${total_bill_amount:,.2f}")
def fill_existing_table(slide, data, keys, slide_title):
    table_shape = find_table_in_slide(slide)

    if not table_shape:
        print(f"No table found in slide '{slide_title}'")
        return False

    table = table_shape.table
    print(f"Found existing table with {len(table.rows)} rows and {len(table.columns)} columns")

    if not data:
        print(f"No data available for '{slide_title}', keeping header only")
        try:
            while len(table.rows) > 1:
                row_to_remove = table.rows[-1]._tr
                table._tbl.remove(row_to_remove)

            if len(table.rows) == 1:
                last_row = table.rows[-1]._tr
                new_row = last_row.__class__(last_row.xml, last_row)
                table._tbl.append(new_row)

            for col_idx in range(min(len(table.columns), len(keys))):
                cell = table.cell(1, col_idx)
                cell.text = "No data available" if col_idx == 0 else ""
        except Exception as e:
            print(f"Error clearing table: {e}")
        return True

    try:
        adjust_table_rows(table, len(data))
    except Exception as e:
        print(f"Error adjusting table rows: {e}")

    if len(table.columns) < len(keys):
        print(f"Warning: Table has {len(table.columns)} columns but need {len(keys)}")
        keys = keys[:len(table.columns)]

    for row_idx, item in enumerate(data, start=1):
        if row_idx >= len(table.rows):
            print(f"Warning: Not enough rows in table for all data (need {len(data) + 1}, have {len(table.rows)})")
            break

        for col_idx, key in enumerate(keys):
            if col_idx >= len(table.columns):
                break

            try:
                text = str(item.get(key, "N/A"))
                cell = table.cell(row_idx, col_idx)
                cell.text = text

                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(11)
                para.alignment = PP_ALIGN.CENTER

                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
            except Exception as e:
                print(f"Error filling cell [{row_idx}, {col_idx}]: {e}")

    print(f"Successfully filled table in '{slide_title}' with {len(data)} data rows")
    return True


def add_optimization_pie_chart(prs, slide_index, categorized_resources):
    slide = prs.slides[slide_index]

    labels = []
    sizes = []
    colors = ["#4CAF50", "#FFC107", "#F44336", "#9E9E9E"]

    for category in ["Optimized", "Under Provisioned", "Over Provisioned", "No Recommendation"]:
        labels.append(category)
        sizes.append(len(categorized_resources.get(category, [])))

    plt.figure(figsize=(4, 4))
    plt.pie(sizes, labels=labels, colors=colors, autopct='%1.1f%%')
    chart_path = "output/resource_distribution_pie.png"
    plt.savefig(chart_path)
    plt.close()

    insert_image_to_slide(slide, chart_path, prs, left=Inches(0.5), top=Inches(1.5))  # Adjust if needed

    # Populate textboxes for each category
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.lower()
            for category in categorized_resources:
                if category.lower() in text:
                    resource_names = categorized_resources[category]
                    shape.text_frame.clear()
                    shape.text_frame.text = f"{category}:\n" + "\n".join(resource_names)


def add_billing_summary_to_slide(prs, slide_index, billing_data):
    slide = prs.slides[slide_index]

    # === Generate Pie Chart ===
    labels = list(billing_data.keys())
    values = list(billing_data.values())

    colors = plt.get_cmap('tab20').colors[:len(labels)]

    fig, ax = plt.subplots()
    wedges, texts, autotexts = ax.pie(
        values, labels=labels, autopct='%1.1f%%',
        startangle=90, colors=colors, textprops={'fontsize': 8}
    )
    ax.axis('equal')

    chart_path = "output/billing_summary_pie.png"
    plt.savefig(chart_path, bbox_inches='tight')
    plt.close()

    # === Insert Pie Chart ===
    chart_left = Inches(5.5)
    chart_top = Inches(1)
    chart_width = Inches(4)
    chart_height = Inches(4)

    slide.shapes.add_picture(chart_path, chart_left, chart_top, width=chart_width, height=chart_height)

    # === Add Service Name and Amount Boxes on the Left ===
    left_x = Inches(0.5)
    top_y = Inches(1)
    spacing = Inches(0.5)

    for idx, (service, amount) in enumerate(billing_data.items()):
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            left_x, top_y + idx * spacing,
            Inches(4.5), Inches(0.4)
        )
        box.text = f"{service}: ${amount:,.2f}"
        text_frame = box.text_frame
        text_frame.paragraphs[0].font.size = Pt(14)
        text_frame.paragraphs[0].font.bold = True
        box.fill.solid()
        box.fill.fore_color.rgb = slide.shapes.title.fill.fore_color.rgb  # match title color if needed
        box.line.color.rgb = box.fill.fore_color.rgb
