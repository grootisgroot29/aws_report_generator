import os
import sys

from pptx import Presentation

from data_collectors.ec2 import get_ec2_instances_with_metrics, get_ec2_backup_metrics
from data_collectors.eks import get_eks_clusters_with_metrics
from data_collectors.iam import get_iam_users_with_metrics
from data_collectors.rds import get_rds_instances_with_metrics, get_rds_backup_metrics
from utils.monthly_metric import get_monthly_metrics
from utils.ppt_edit import add_billing_summary_to_slide
from utils.monthly_billing import get_monthly_billing_data
from utils.plots import plot_metrics
from utils.ppt_edit import (
    find_slide_by_title,
    fill_existing_table,
)
from utils.ppt_edit import insert_image_to_slide, update_textbox_with_resource_name
from utils.ppt_edit import update_resource_counts_on_slide


def main(path="template/Report1.pptx", output_filename="output/AWS_Services_Report.pptx"):
    print(f"Loading template from: {path}")
    try:
        prs = Presentation(path)
        print("Template loaded successfully")
    except Exception as e:
        print(f"Error loading template: {e}")
        return

    print("\nCollecting AWS metrics...")

    print("\nCollecting EC2 data...")
    ec2_data = get_ec2_instances_with_metrics()
    print(f"Found {len(ec2_data)} EC2 instances")

    print("\nCollecting RDS data...")
    rds_data = get_rds_instances_with_metrics()
    print(f"Found {len(rds_data)} RDS instances")
    # Slide index for EC2 & RDS count summary (e.g., slide 2 → index 1)
    update_resource_counts_on_slide(
        prs,
        slide_index=1,  # Assuming first slide
        ec2_count=len(ec2_data),
        rds_count=len(rds_data),
        # total_bill_amount=total_amount  # Float value from billing
    )

    #After EC2/RDS data is collected
    print("\nFetching optimization classifications...")
    categorized_resources = get_aws_optimization_status()

    print("\nAdding resource classification chart to slide 9...")
    add_optimization_pie_chart(prs, slide_index=8, categorized_resources=categorized_resources)

    print("\nCollecting EKS data...")
    eks_data = get_eks_clusters_with_metrics()
    print(f"Found {len(eks_data)} EKS clusters")

    print("\nCollecting IAM data...")
    iam_data = get_iam_users_with_metrics()
    print(f"Found {len(iam_data)} IAM users")

    print("\nCollecting EC2 backup data...")
    ec2_backup_data = get_ec2_backup_metrics()
    print(f"Found backup info for {len(ec2_backup_data)} EC2 instances")

    print("\nCollecting RDS backup data...")
    rds_backup_data = get_rds_backup_metrics()
    print(f"Found backup info for {len(rds_backup_data)} RDS instances")

    output_dir = "output"
    os.makedirs(output_dir, exist_ok=True)


    ec2_graph_ready = []
    for ec2 in ec2_data:
        #Check that the metrics are present and not "N/A" or "Error"
        if all(
            ec2.get(metric_key) not in ["N/A", "Error", None]
            for metric_key in ['monthly_cpu_usage_(%)', 'monthly_memory_usage_(%)', 'monthly_disk_usage_(%)']
        ):
            ec2_graph_ready.append(ec2)

    #Filter RDS instances with all required metrics for graphs
    rds_graph_ready = []
    for rds in rds_data:
        metrics = get_monthly_metrics(
            resource_id=rds['db_identifier'],
            metric_names=["CPUUtilization", "FreeableMemory"],
            namespace="AWS/RDS",
            dimension_name="DBInstanceIdentifier"
        )
        if metrics and all(metrics.get(m) for m in ["CPUUtilization", "FreeableMemory"]):
            rds_graph_ready.append(rds)

    def add_ec2_graphs_to_slide(prs, slide_index, ec2_instances, output_dir):
        slide = prs.slides[slide_index]

        for ec2 in ec2_instances:
            instance_id = ec2['instance_id']  # Use instance_id key
            server_name = ec2['server_name']

            print(f"Generating graph for: {server_name}")

            cpu_metrics = get_monthly_metrics(
                resource_id=instance_id,
                metric_names=["CPUUtilization"],
                namespace="AWS/EC2",
                dimension_name="InstanceId"
            ) or {}

            cwagent_metrics = get_monthly_metrics(
                resource_id=instance_id,
                metric_names=["mem_used_percent", "disk_used_percent"],
                namespace="CWAgent",
                dimension_name="InstanceId"
            ) or {}

            metrics = {**cpu_metrics, **cwagent_metrics}

            if any(metrics.values()):
                graph_path = os.path.join(output_dir, f"{server_name}_ec2_metrics.png")
                plot_metrics(metrics, f"EC2 Metrics for {server_name}", graph_path)

                if os.path.exists(graph_path):
                    print(f"Graph generated at: {graph_path}")
                    insert_image_to_slide(slide, graph_path, prs)
                    update_textbox_with_resource_name(slide, "server name:", server_name)
                else:
                    print(f"Graph file not found: {graph_path}")
            else:
                print(f"No available metrics to plot for: {server_name}")

    def add_rds_graphs_to_slide(prs, slide_index, rds_instances):
        slide = prs.slides[slide_index]
        for rds in rds_instances:
            metrics = get_monthly_metrics(
                resource_id=rds['db_identifier'],
                metric_names=["CPUUtilization", "FreeableMemory"],
                namespace="AWS/RDS",
                dimension_name="DBInstanceIdentifier"
            )
            if metrics:
                graph_path = os.path.join(output_dir, f"{rds['db_identifier']}_rds_metrics.png")
                plot_metrics(metrics, f"RDS Metrics for {rds['db_identifier']}", graph_path)
                insert_image_to_slide(slide, graph_path, prs)
                update_textbox_with_resource_name(slide, "RDS Name:", rds['db_identifier'])

    print("\nAdding EC2 graphs to slide 5...")
    add_ec2_graphs_to_slide(prs, slide_index=4, ec2_instances=ec2_graph_ready,output_dir=output_dir)  # zero-based index

    print("\nAdding RDS graphs to slide 8...")
    add_rds_graphs_to_slide(prs, slide_index=7, rds_instances=rds_graph_ready)

    slide_data_map = {
        'EC2 Instances': {
            'keys': ['server_name', 'specification', 'status', 'monthly_cpu_usage_(%)', 'monthly_memory_usage_(%)',
                     'monthly_disk_usage_(%)'],
            'data': ec2_data
        },
        'Relational Databases': {
            'keys': ['db_identifier', 'database_name', 'engine', 'status', 'storage_used/allocated',
                     'monthly_cpu_usage_(%)'],
            'data': rds_data
        },
        'EKS Clusters': {
            'keys': ['cluster_name', 'kubernetes_version', 'support_period', 'node_groups'],
            'data': eks_data
        },
        'IAM Users': {
            'keys': ['username', 'groups', 'mfa', 'active_key_age_(days)'],
            'data': iam_data
        },
        'EC2 Backups': {
            'keys': ['instance_name', 'ami_id', 'backup_date', 'status', 'retention_days', 'next_backup'],
            'data': ec2_backup_data
        },
        'RDS Backup': {
            'keys': ['db_name', 'snapshot_id', 'date', 'status', 'retention_days', 'size_gb'],
            'data': rds_backup_data
        }
    }

    print(f"\nUpdating PowerPoint slides...")
    for slide_title, config in slide_data_map.items():
        print(f"\nLooking for slide: '{slide_title}'")
        slide = find_slide_by_title(prs, slide_title)
        if slide:
            print(f"Updating slide: '{slide_title}'")
            try:
                fill_existing_table(slide, config['data'], config['keys'], slide_title)
            except Exception as e:
                print(f"Error updating slide '{slide_title}': {e}")
        else:
            print(f"Could not find slide with title: '{slide_title}'")
    print("\nCollecting billing data...")
    billing_data = get_monthly_billing_data()
    print("Adding billing summary to slide 19...")
    add_billing_summary_to_slide(prs, slide_index=18, billing_data=billing_data)

    print(f"\nSaving presentation...")
    try:
        prs.save(output_filename)
        print(f"Report saved as: {output_filename}")
        print(f"Report generation completed successfully!")
    except Exception as e:
        print(f"Error saving presentation: {e}")
        try:
            fallback_filename = f"output/fallback_{output_filename}"
            prs.save(fallback_filename)
            print(f"Report saved as: {fallback_filename}")
        except Exception as e2:
            print(f"Failed to save presentation: {e2}")


if __name__ == "__main__":
    if len(sys.argv) > 1:
        template_path = sys.argv[1]
        main(path=template_path)
    else:
        main()
